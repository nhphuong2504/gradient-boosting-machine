"""Build a PowerPoint version of the Friedman TreeBoost presentation.

Run from the project root::

    venv\\Scripts\\python.exe slides\\build_pptx.py

The script writes ``slides/presentation.pptx`` mirroring the reveal.js deck.
Mathematical formulas are rendered through matplotlib's ``mathtext`` engine
into transparent PNGs (cached in ``slides/_math_cache``) and embedded as
pictures, so every formula appears as proper typeset notation -- not pseudo-
code -- and looks the same in any version of PowerPoint or Keynote without
needing a LaTeX install.
"""

from __future__ import annotations

import hashlib
from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib import rcParams
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from lxml import etree

# Computer Modern fonts -> the math images look like the symbols in Friedman's paper.
rcParams["mathtext.fontset"] = "cm"


# ---------------------------------------------------------------------------
# Theme
# ---------------------------------------------------------------------------

ACCENT = RGBColor(0x1F, 0x4E, 0x79)
ACCENT_SOFT = RGBColor(0xE8, 0xEE, 0xF5)
TEXT = RGBColor(0x1F, 0x29, 0x37)
MUTED = RGBColor(0x6B, 0x72, 0x80)
GOOD_BG = RGBColor(0xDC, 0xFC, 0xE7)
GOOD_TEXT = RGBColor(0x16, 0x65, 0x34)
WARN_BG = RGBColor(0xFE, 0xF3, 0xC7)
WARN_TEXT = RGBColor(0xB5, 0x47, 0x08)
CARD_BG = RGBColor(0xFA, 0xFA, 0xFA)
BORDER = RGBColor(0xE5, 0xE7, 0xEB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT_BODY = "Calibri"
FONT_MONO = "Consolas"

ROOT = Path(__file__).resolve().parent.parent
RESULTS = ROOT / "experiments" / "results"
OUTPUT = ROOT / "slides" / "presentation.pptx"
MATH_CACHE = ROOT / "slides" / "_math_cache"
MATH_CACHE.mkdir(parents=True, exist_ok=True)


# ---------------------------------------------------------------------------
# Math rendering
# ---------------------------------------------------------------------------
#
# matplotlib's mathtext is a self-contained subset of LaTeX -- it supports
# \frac, \sum, \arg\min, \mathbb, \partial, \tilde, etc., but it does NOT
# support multi-line environments such as ``cases``, ``align`` or ``array``.
# We therefore render each *single-line* formula via mathtext, and stack
# multiple PNGs vertically when we need a piecewise definition or a multi-
# line algorithm display. Each rendered formula is cached on disk so reruns
# of this script are fast.

# A consistent dark text colour across all rendered math (matches body TEXT).
_MATH_HEX = "#1f2937"


def _math_cache_key(latex: str, fontsize: int, color: str) -> str:
    """Stable filename for a (latex, fontsize, color) triple."""
    h = hashlib.sha1(f"{latex}|{fontsize}|{color}".encode("utf-8")).hexdigest()
    return h[:16]


def render_math(latex: str, *, fontsize: int = 28,
                color: str = _MATH_HEX) -> Path:
    """Render a single-line LaTeX math expression to a transparent PNG.

    Returns the path to the cached image file.
    """
    key = _math_cache_key(latex, fontsize, color)
    out = MATH_CACHE / f"{key}.png"
    if out.exists():
        return out
    fig, ax = plt.subplots(figsize=(0.01, 0.01))
    ax.set_axis_off()
    fig.patch.set_alpha(0)
    ax.patch.set_alpha(0)
    ax.text(
        0.5, 0.5, f"${latex}$",
        fontsize=fontsize, color=color,
        ha="center", va="center",
    )
    fig.savefig(
        out,
        dpi=300,
        transparent=True,
        bbox_inches="tight",
        pad_inches=0.02,
    )
    plt.close(fig)
    return out


def _math_image_size(path: Path):
    with Image.open(path) as img:
        return img.size  # (w_pixels, h_pixels)


def add_math_image(slide, x, y, max_w, max_h, latex, *,
                   fontsize: int = 28, color: str = _MATH_HEX,
                   align: str = "center"):
    """Insert a rendered math PNG fitted within (max_w, max_h).

    Aspect ratio is preserved. ``align`` controls horizontal placement
    when the rendered width is smaller than ``max_w``: ``"center"`` (default)
    or ``"left"``.
    """
    path = render_math(latex, fontsize=fontsize, color=color)
    iw, ih = _math_image_size(path)
    aspect = iw / ih
    target_w = max_w
    target_h = int(target_w / aspect)
    if target_h > max_h:
        target_h = max_h
        target_w = int(target_h * aspect)
    if align == "center":
        cx = x + (max_w - target_w) // 2
    else:  # "left"
        cx = x
    cy = y + (max_h - target_h) // 2
    return slide.shapes.add_picture(
        str(path), cx, cy, width=target_w, height=target_h,
    )


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def add_blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank layout


def fill(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def no_line(shape):
    shape.line.fill.background()


def set_line(shape, color: RGBColor, width_pt: float = 1.0):
    shape.line.color.rgb = color
    shape.line.width = Pt(width_pt)


def add_rect(slide, x, y, w, h, fill_color=None, line_color=None,
             shape_type=MSO_SHAPE.RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, x, y, w, h)
    if fill_color is not None:
        fill(shape, fill_color)
    else:
        shape.fill.background()
    if line_color is None:
        no_line(shape)
    else:
        set_line(shape, line_color, 0.75)
    shape.shadow.inherit = False
    return shape


def add_textbox(slide, x, y, w, h, text="", *, font=FONT_BODY, size=18,
                bold=False, italic=False, color=TEXT, align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if text:
        run = p.add_run()
        run.text = text
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def add_paragraph(text_frame, text="", *, font=FONT_BODY, size=16, bold=False,
                  italic=False, color=TEXT, align=PP_ALIGN.LEFT,
                  bullet=False, level=0, space_before=2, space_after=2):
    p = text_frame.add_paragraph()
    p.alignment = align
    p.level = level
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    if bullet:
        # Add a manual bullet dot — keeps formatting consistent across themes.
        run = p.add_run()
        run.text = "• "
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
    if text:
        run = p.add_run()
        run.text = text
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return p


def add_run(paragraph, text, *, font=FONT_BODY, size=16, bold=False,
            italic=False, color=TEXT):
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return run


def add_first_run(text_frame, text, *, font=FONT_BODY, size=16, bold=False,
                  italic=False, color=TEXT, align=PP_ALIGN.LEFT):
    p = text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return run


def add_title(slide, text, *, top=Inches(0.30), left=Inches(0.55),
              width=Inches(12.2), height=Inches(0.7)):
    tb = add_textbox(slide, left, top, width, height, text=text,
                     size=30, bold=True, color=ACCENT, anchor=MSO_ANCHOR.TOP)
    return tb


def add_subtitle(slide, text, *, top=Inches(0.95), left=Inches(0.55),
                 width=Inches(12.2), height=Inches(0.4)):
    tb = add_textbox(slide, left, top, width, height, text=text,
                     size=15, color=MUTED, italic=True)
    return tb


def add_accent_underline(slide, *, top=Inches(0.95), left=Inches(0.55),
                         width=Inches(0.8)):
    bar = add_rect(slide, left, top, width, Emu(45720),  # ~0.05"
                   fill_color=ACCENT)
    return bar


def add_footer(slide, page, total, *, paper="Friedman (2001) — Greedy Function Approximation"):
    add_textbox(slide, Inches(0.55), Inches(7.05), Inches(8.0), Inches(0.3),
                text=paper, size=10, color=MUTED, italic=True)
    add_textbox(slide, Inches(11.0), Inches(7.05), Inches(2.0), Inches(0.3),
                text=f"{page} / {total}", size=10, color=MUTED,
                align=PP_ALIGN.RIGHT)


def add_keypoint_box(slide, x, y, w, h, text, *, size=14):
    """Accent-soft callout box, vertically centered text."""
    rect = add_rect(slide, x, y, w, h, fill_color=ACCENT_SOFT)
    # Left accent bar
    add_rect(slide, x, y, Inches(0.06), h, fill_color=ACCENT)
    tb = slide.shapes.add_textbox(x + Inches(0.18), y, w - Inches(0.25), h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.name = FONT_BODY
    run.font.size = Pt(size)
    run.font.color.rgb = TEXT
    return rect


def add_keypoint_math(slide, x, y, w, h, latex, *, fontsize=36):
    """Accent-soft callout containing a centred rendered math image."""
    rect = add_rect(slide, x, y, w, h, fill_color=ACCENT_SOFT)
    add_rect(slide, x, y, Inches(0.06), h, fill_color=ACCENT)
    inner_pad_x = Inches(0.30)
    inner_pad_y = Inches(0.12)
    add_math_image(
        slide,
        x + inner_pad_x,
        y + inner_pad_y,
        w - 2 * inner_pad_x,
        h - 2 * inner_pad_y,
        latex,
        fontsize=fontsize,
    )
    return rect


def add_stacked_math(slide, x, y, w, line_height, latex_lines, *,
                     fontsize: int = 22, gap=Emu(45720),
                     align: str = "left", color: str = _MATH_HEX):
    """Render a list of math expressions and stack them vertically.

    Used for piecewise definitions and the Algorithm-1 display where
    matplotlib's mathtext lacks a multi-line environment.
    """
    cur_y = y
    for latex in latex_lines:
        path = render_math(latex, fontsize=fontsize, color=color)
        iw, ih = _math_image_size(path)
        aspect = iw / ih
        target_h = line_height
        target_w = int(target_h * aspect)
        if target_w > w:
            target_w = w
            target_h = int(target_w / aspect)
        if align == "center":
            cx = x + (w - target_w) // 2
        else:  # left
            cx = x
        slide.shapes.add_picture(
            str(path), cx, cur_y, width=target_w, height=target_h,
        )
        cur_y += line_height + gap
    return cur_y


def add_pill(slide, x, y, text, *, fg=ACCENT, bg=ACCENT_SOFT, size=11):
    width = Inches(0.05 + 0.085 * len(text))
    height = Inches(0.28)
    pill = add_rect(slide, x, y, width, height, fill_color=bg,
                    shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    pill.adjustments[0] = 0.5
    tf = pill.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text.upper()
    run.font.name = FONT_BODY
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = fg
    return pill, width


def add_code_block(slide, x, y, w, h, code, *, size=12):
    rect = add_rect(slide, x, y, w, h, fill_color=CARD_BG, line_color=BORDER)
    tb = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.1),
                                  w - Inches(0.3), h - Inches(0.2))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    lines = code.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        run = p.add_run()
        run.text = line if line else "\u00A0"
        run.font.name = FONT_MONO
        run.font.size = Pt(size)
        run.font.color.rgb = TEXT
    return rect


def add_picture_card(slide, x, y, w, h, image_path: Path):
    rect = add_rect(slide, x, y, w, h, fill_color=WHITE, line_color=BORDER)
    if not image_path.exists():
        add_textbox(slide, x, y, w, h, text=f"[missing] {image_path.name}",
                    size=14, color=WARN_TEXT, align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE)
        return rect
    pad = Inches(0.1)
    pic = slide.shapes.add_picture(
        str(image_path), x + pad, y + pad,
        width=w - 2 * pad, height=h - 2 * pad,
    )
    return pic


def add_table(slide, x, y, w, h, data, *, header_fill=ACCENT_SOFT,
              header_color=ACCENT, font_size=14, num_cols_right=None):
    rows = len(data)
    cols = len(data[0])
    table_shape = slide.shapes.add_table(rows, cols, x, y, w, h)
    table = table_shape.table
    num_cols_right = num_cols_right or set()
    for ri, row in enumerate(data):
        for ci, value in enumerate(row):
            cell = table.cell(ri, ci)
            cell.margin_left = Emu(45720)
            cell.margin_right = Emu(45720)
            cell.margin_top = Emu(22860)
            cell.margin_bottom = Emu(22860)
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_fill if ri == 0 else WHITE
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = (PP_ALIGN.RIGHT if ci in num_cols_right and ri > 0
                           else PP_ALIGN.LEFT)
            run = p.add_run()
            run.text = str(value)
            run.font.name = FONT_BODY
            run.font.size = Pt(font_size)
            run.font.bold = (ri == 0)
            run.font.color.rgb = header_color if ri == 0 else TEXT
    return table


# ---------------------------------------------------------------------------
# Build the deck
# ---------------------------------------------------------------------------

TOTAL_SLIDES = 20

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


# ---------- Slide 1 — Title ----------------------------------------------
def slide_1_title():
    slide = add_blank_slide(prs)
    # Banner
    add_rect(slide, 0, 0, prs.slide_width, Inches(7.5), fill_color=WHITE)
    add_rect(slide, 0, Inches(2.4), prs.slide_width, Inches(0.06),
             fill_color=ACCENT)
    add_rect(slide, 0, Inches(5.0), prs.slide_width, Inches(0.06),
             fill_color=ACCENT_SOFT)

    add_textbox(slide, Inches(0.7), Inches(2.7), Inches(12), Inches(0.5),
                text="Greedy Function Approximation:",
                size=42, bold=True, color=ACCENT)
    add_textbox(slide, Inches(0.7), Inches(3.35), Inches(12), Inches(0.7),
                text="A Gradient Boosting Machine",
                size=42, bold=True, color=ACCENT)

    add_textbox(slide, Inches(0.7), Inches(4.25), Inches(12), Inches(0.4),
                text="Jerome H. Friedman   ·   Annals of Statistics, 2001",
                size=18, color=TEXT)

    add_textbox(slide, Inches(0.7), Inches(5.25), Inches(12), Inches(0.35),
                text="Reading + from-scratch implementation",
                size=15, color=MUTED, italic=True)
    add_textbox(slide, Inches(0.7), Inches(5.65), Inches(12), Inches(0.35),
                text="LS Boost   ·   LAD TreeBoost   ·   M (Huber) TreeBoost",
                size=14, color=MUTED, italic=True)


# ---------- Slide 2 — Roadmap --------------------------------------------
def slide_2_roadmap():
    slide = add_blank_slide(prs)
    add_title(slide, "Roadmap")
    add_accent_underline(slide)

    items = [
        ("1.", "Why we need this paper — the problem and the literature"),
        ("2.", "The key insight: gradient descent in function space"),
        ("3.", "Three regression algorithms: LS, LAD, Huber TreeBoost"),
        ("4.", "The terminal-region update — the practical heart of TreeBoost"),
        ("5.", "From-scratch implementation in Python"),
        ("6.", "Experiments: convergence, robustness, shrinkage"),
        ("7.", "What the results reveal & the paper's lasting impact"),
    ]
    y = Inches(1.55)
    for num, text in items:
        add_textbox(slide, Inches(0.9), y, Inches(0.5), Inches(0.55),
                    text=num, size=22, bold=True, color=ACCENT)
        add_textbox(slide, Inches(1.4), y, Inches(11), Inches(0.55),
                    text=text, size=20, color=TEXT)
        y += Inches(0.65)
    add_footer(slide, 2, TOTAL_SLIDES)


# ---------- Slide 3 — Problem --------------------------------------------
def slide_3_problem():
    slide = add_blank_slide(prs)
    add_title(slide, "The function-estimation problem")
    add_accent_underline(slide)

    add_textbox(
        slide, Inches(0.55), Inches(1.40), Inches(12.2), Inches(0.45),
        text=(
            "Given training data {(x_i, y_i)} for i = 1, …, N, find a "
            "function F : X → ℝ minimizing the expected loss"
        ),
        size=18, color=TEXT,
    )
    add_keypoint_math(
        slide, Inches(2.0), Inches(2.00), Inches(9.3), Inches(1.0),
        latex=r"F^{*} \;=\; \arg\min_{F}\; \mathbb{E}_{y,\,\mathbf{x}}\,"
              r"L\!\left(y,\, F(\mathbf{x})\right)",
        fontsize=34,
    )

    add_textbox(slide, Inches(0.55), Inches(3.20), Inches(12), Inches(0.4),
                text="Concrete losses:", size=16, bold=True, color=ACCENT)
    # Each row: bold label on the left, rendered math on the right.
    rows = [
        ("Regression:",
         r"L(y, F) \;=\; \frac{1}{2}(y - F)^{2}"
         r"\;\;\text{or}\;\;"
         r"L(y, F) \;=\; |\,y - F\,|"),
        ("Classification:",
         r"L(y, F) \;=\; \log\!\left(1 + e^{-2\,y\,F}\right)"
         r",\qquad y \in \{-1,\, +1\}"),
    ]
    y = Inches(3.62)
    for label, latex in rows:
        add_textbox(slide, Inches(1.0), y, Inches(2.2), Inches(0.5),
                    text="•  " + label, size=17, bold=True, color=ACCENT,
                    anchor=MSO_ANCHOR.MIDDLE)
        add_math_image(
            slide, Inches(3.2), y, Inches(9.5), Inches(0.5),
            latex=latex, fontsize=20, align="left",
        )
        y += Inches(0.55)

    add_keypoint_box(
        slide, Inches(0.55), Inches(5.55), Inches(12.2), Inches(1.15),
        "Most existing methods restrict F to a parametric family and tackle a "
        "hard joint optimization. This paper takes a non-parametric, stagewise route.",
        size=16,
    )
    add_footer(slide, 3, TOTAL_SLIDES)


# ---------- Slide 4 — Position in literature -----------------------------
def slide_4_literature():
    slide = add_blank_slide(prs)
    add_title(slide, "Where this paper sits")
    add_accent_underline(slide)

    col_top = Inches(1.45)
    col_h = Inches(4.6)

    # Left column — Before
    add_rect(slide, Inches(0.55), col_top, Inches(6.05), col_h,
             fill_color=WHITE, line_color=BORDER)
    add_textbox(slide, Inches(0.75), Inches(1.55), Inches(5.5), Inches(0.45),
                text="Before (1995–2000)", size=20, bold=True, color=ACCENT)
    before = [
        "AdaBoost (Freund & Schapire, 1996) — re-weights examples; tied to "
        "exponential loss.",
        "LogitBoost (Friedman, Hastie, Tibshirani, 2000) — binomial likelihood "
        "via Newton steps.",
        "Basis-function methods: MARS, RBF networks, SVM, neural nets.",
    ]
    y = Inches(2.10)
    for item in before:
        add_textbox(slide, Inches(0.75), y, Inches(0.3), Inches(0.4),
                    text="•", size=18, bold=True, color=ACCENT)
        add_textbox(slide, Inches(1.05), y, Inches(5.4), Inches(1.2),
                    text=item, size=14, color=TEXT)
        y += Inches(0.95)

    # Right column — This paper
    add_rect(slide, Inches(6.75), col_top, Inches(6.05), col_h,
             fill_color=ACCENT_SOFT, line_color=BORDER)
    add_textbox(slide, Inches(6.95), Inches(1.55), Inches(5.5), Inches(0.45),
                text="This paper", size=20, bold=True, color=ACCENT)
    this = [
        "One unifying view: boosting is gradient descent in function space.",
        "Plug in any differentiable loss.",
        "Concrete algorithms for LS, LAD, Huber, and K-class logistic.",
        "Adds tools to interpret tree boosters.",
    ]
    y = Inches(2.10)
    for item in this:
        add_textbox(slide, Inches(6.95), y, Inches(0.3), Inches(0.4),
                    text="•", size=18, bold=True, color=ACCENT)
        add_textbox(slide, Inches(7.25), y, Inches(5.4), Inches(0.85),
                    text=item, size=14, color=TEXT)
        y += Inches(0.65)

    add_textbox(
        slide, Inches(0.55), Inches(6.35), Inches(12.2), Inches(0.5),
        text=("After: XGBoost (2014), LightGBM (2017), CatBoost (2017) — "
              "all extend this paradigm."),
        size=13, italic=True, color=MUTED,
    )
    add_footer(slide, 4, TOTAL_SLIDES)


# ---------- Slide 5 — Stagewise additive expansions ----------------------
def slide_5_additive():
    slide = add_blank_slide(prs)
    add_title(slide, "Stagewise additive expansions")
    add_accent_underline(slide)

    add_textbox(slide, Inches(0.55), Inches(1.45), Inches(12), Inches(0.45),
                text="Approximate F by a sum of simple \u201cweak learners\u201d:",
                size=18, color=TEXT)
    add_keypoint_math(
        slide, Inches(2.0), Inches(2.00), Inches(9.3), Inches(1.0),
        latex=r"F\!\left(\mathbf{x};\,\{\beta_m,\,\mathbf{a}_m\}_{m=1}^{M}\right)"
              r" \;=\; \sum_{m=1}^{M}\, \beta_m\, h(\mathbf{x};\,\mathbf{a}_m)",
        fontsize=32,
    )

    bullets = [
        "h(x ; a) — a small parametric function (e.g. a regression tree).",
        "Joint optimization over all (β_m, a_m) is generally hard.",
    ]
    y = Inches(3.20)
    for item in bullets:
        add_textbox(slide, Inches(0.9), y, Inches(0.3), Inches(0.4),
                    text="•", size=18, bold=True, color=ACCENT)
        add_textbox(slide, Inches(1.2), y, Inches(11), Inches(0.4),
                    text=item, size=16, color=TEXT)
        y += Inches(0.45)

    add_textbox(slide, Inches(0.55), Inches(4.20), Inches(12), Inches(0.4),
                text="Greedy stagewise alternative — at each step m:",
                size=17, bold=True, color=ACCENT)
    add_keypoint_math(
        slide, Inches(0.55), Inches(4.75), Inches(12.2), Inches(1.20),
        latex=r"(\beta_m,\, \mathbf{a}_m) \;=\; \arg\min_{\beta,\,\mathbf{a}}"
              r" \sum_{i=1}^{N}\, L\!\left(y_i,\; F_{m-1}(\mathbf{x}_i)"
              r" + \beta\, h(\mathbf{x}_i;\, \mathbf{a})\right)",
        fontsize=26,
    )

    add_textbox(slide, Inches(0.55), Inches(6.10), Inches(12), Inches(0.4),
                text="Eq. (9) of the paper. Previous terms are frozen — no re-fitting.",
                size=12, italic=True, color=MUTED)
    add_footer(slide, 5, TOTAL_SLIDES)


# ---------- Slide 6 — Key insight ----------------------------------------
def slide_6_insight():
    slide = add_blank_slide(prs)
    add_title(slide, "The key insight")
    add_subtitle(slide, "Boosting = numerical optimization in function space")
    add_accent_underline(slide, top=Inches(1.40))

    add_textbox(slide, Inches(0.55), Inches(1.65), Inches(12), Inches(0.5),
                text=("Treat F(x_i) at each training point as a free parameter. "
                      "The negative gradient at iteration m is"),
                size=16, color=TEXT)

    add_keypoint_math(
        slide, Inches(0.95), Inches(2.20), Inches(11.4), Inches(1.30),
        latex=r"\tilde{y}_i \;=\; -\!\left[\,\frac{\partial L\!\left(y_i,\,"
              r"F(\mathbf{x}_i)\right)}{\partial F(\mathbf{x}_i)}\,\right]_{"
              r"F(\mathbf{x})\,=\,F_{m-1}(\mathbf{x})}",
        fontsize=30,
    )

    add_textbox(slide, Inches(0.55), Inches(3.65), Inches(12), Inches(0.5),
                text=("These are the pseudo-responses: the direction in which "
                      "moving F(x_i) most reduces the loss."),
                size=16, color=TEXT)

    add_keypoint_box(
        slide, Inches(0.55), Inches(4.55), Inches(12.2), Inches(2.0),
        "Why this matters:  the gradient is only defined at the training "
        "points. Fit a weak learner h(x ; a_m) to the pseudo-responses by "
        "least squares — that's the smooth, generalizable approximation to "
        "the unconstrained gradient.",
        size=16,
    )
    add_footer(slide, 6, TOTAL_SLIDES)


# ---------- Slide 7 — Algorithm 1 ----------------------------------------
def slide_7_algo1():
    slide = add_blank_slide(prs)
    add_title(slide, "Algorithm 1 — Gradient Boost (generic)")
    add_accent_underline(slide)

    # Light card behind the algorithm display.
    algo_box_x = Inches(0.55)
    algo_box_y = Inches(1.45)
    algo_box_w = Inches(12.2)
    algo_box_h = Inches(2.85)
    add_rect(slide, algo_box_x, algo_box_y, algo_box_w, algo_box_h,
             fill_color=CARD_BG, line_color=BORDER)

    # Each algorithm step: (left label, rendered math, optional comment).
    steps = [
        ("1.",
         r"F_{0}(\mathbf{x}) \;=\; \arg\min_{\rho}\,"
         r"\sum_{i=1}^{N}\, L(y_i,\, \rho)",
         "initialize"),
        ("2.",
         r"\mathbf{for}\;\; m = 1, 2, \ldots, M\!:",
         ""),
        ("3.",
         r"\tilde{y}_i \;=\; -\!\left[\,\frac{\partial L(y_i,\, F(\mathbf{x}_i))}"
         r"{\partial F(\mathbf{x}_i)}\,\right]_{F = F_{m-1}}",
         "pseudo-responses"),
        ("4.",
         r"\mathbf{a}_m \;=\; \arg\min_{\mathbf{a},\,\beta}\,"
         r"\sum_{i=1}^{N}\, \left[\,\tilde{y}_i - \beta\, h(\mathbf{x}_i;\,\mathbf{a})\,\right]^{2}",
         "least-squares fit"),
        ("5.",
         r"\rho_m \;=\; \arg\min_{\rho}\,"
         r"\sum_{i=1}^{N}\, L\!\left(y_i,\, F_{m-1}(\mathbf{x}_i) + \rho\, h(\mathbf{x}_i;\,\mathbf{a}_m)\right)",
         "1-D line search"),
        ("6.",
         r"F_m(\mathbf{x}) \;=\; F_{m-1}(\mathbf{x}) \,+\, \rho_m\, h(\mathbf{x};\,\mathbf{a}_m)",
         "update"),
    ]

    line_height = Inches(0.40)
    line_gap = Emu(34290)  # ~0.038"
    cur_y = algo_box_y + Inches(0.18)
    label_x = algo_box_x + Inches(0.30)
    math_x = algo_box_x + Inches(0.85)
    math_w = Inches(8.6)
    comment_x = algo_box_x + Inches(9.65)

    for label, latex, comment in steps:
        # Step number
        add_textbox(
            slide, label_x, cur_y, Inches(0.5), line_height,
            text=label, size=14, bold=True, color=ACCENT,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        # Indent inner-loop steps so the structure reads
        x_offset = Inches(0.45) if label in ("3.", "4.", "5.", "6.") else Inches(0)
        # Rendered math
        add_math_image(
            slide,
            math_x + x_offset,
            cur_y,
            math_w - x_offset,
            line_height,
            latex=latex,
            fontsize=18,
            align="left",
        )
        # Right-side comment
        if comment:
            add_textbox(
                slide, comment_x, cur_y, Inches(2.7), line_height,
                text="# " + comment, size=12, italic=True, color=MUTED,
                font=FONT_MONO, anchor=MSO_ANCHOR.MIDDLE,
            )
        cur_y += line_height + line_gap

    notes = [
        ("Line 3", " is loss-specific. Everything else is the same machinery."),
        ("Line 4", " is least squares — fast and well-understood, regardless of the original loss."),
        ("Line 5", " is a 1-D line search in the original loss."),
    ]
    y = Inches(4.55)
    for label, body in notes:
        add_textbox(slide, Inches(0.9), y, Inches(0.3), Inches(0.4),
                    text="•", size=18, bold=True, color=ACCENT)
        tb = slide.shapes.add_textbox(Inches(1.2), y, Inches(11.5), Inches(0.55))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = label
        run.font.name = FONT_BODY
        run.font.size = Pt(16)
        run.font.bold = True
        run.font.color.rgb = ACCENT
        run = p.add_run()
        run.text = body
        run.font.name = FONT_BODY
        run.font.size = Pt(16)
        run.font.color.rgb = TEXT
        y += Inches(0.6)
    add_footer(slide, 7, TOTAL_SLIDES)


# ---------- Slide 8 — LS Boost -------------------------------------------
def slide_8_ls():
    slide = add_blank_slide(prs)
    add_title(slide, "Algorithm 2 — LS Boost")
    add_accent_underline(slide)

    # "With L(y, F) = ½(y - F)²:" rendered as math.
    add_textbox(slide, Inches(0.55), Inches(1.45), Inches(2.5), Inches(0.5),
                text="With", size=18, color=TEXT, anchor=MSO_ANCHOR.MIDDLE)
    add_math_image(
        slide, Inches(1.20), Inches(1.45), Inches(5.5), Inches(0.5),
        latex=r"L(y, F) \;=\; \frac{1}{2}\,(y - F)^{2}\, :",
        fontsize=22, align="left",
    )

    # Three bullets: bold label + rendered math expression (or prose).
    # is_math=False marks rows whose value should be rendered as plain text.
    rows = [
        ("Pseudo-response:",
         r"\tilde{y}_i \;=\; y_i - F_{m-1}(\mathbf{x}_i)"
         r"\qquad\text{(residuals).}",
         True),
        ("Line 4:",
         "fits the current residuals — the line search becomes trivial.",
         False),
        ("Initial constant:",
         r"F_0 \;=\; \bar{y}.",
         True),
    ]
    y = Inches(2.10)
    for label, content, is_math in rows:
        add_textbox(slide, Inches(0.9), y, Inches(0.3), Inches(0.5),
                    text="•", size=18, bold=True, color=ACCENT,
                    anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(1.2), y, Inches(2.7), Inches(0.5),
                    text=label, size=16, bold=True, color=ACCENT,
                    anchor=MSO_ANCHOR.MIDDLE)
        if is_math:
            add_math_image(
                slide, Inches(3.95), y, Inches(8.7), Inches(0.5),
                latex=content, fontsize=18, align="left",
            )
        else:
            add_textbox(
                slide, Inches(3.95), y, Inches(8.7), Inches(0.5),
                text=content, size=15, color=TEXT,
                anchor=MSO_ANCHOR.MIDDLE,
            )
        y += Inches(0.55)

    add_keypoint_box(
        slide, Inches(0.55), Inches(3.7), Inches(12.2), Inches(0.85),
        "Reduces to the classical \u201cfit residuals iteratively\u201d recipe — "
        "the reality check for the framework.",
        size=15,
    )

    code = (
        "# src/treeboost/losses.py\n"
        "class LeastSquaresLoss(Loss):\n"
        "    def initial_prediction(self, y):     return float(np.mean(y))\n"
        "    def negative_gradient(self, y, F):   return y - F            # residuals\n"
        "    def leaf_update(self, y, F, idx):    return float(np.mean((y - F)[idx]))"
    )
    add_code_block(slide, Inches(0.55), Inches(4.75),
                   Inches(12.2), Inches(2.0), code, size=13)
    add_footer(slide, 8, TOTAL_SLIDES)


# ---------- Slide 9 — LAD ------------------------------------------------
def slide_9_lad():
    slide = add_blank_slide(prs)
    add_title(slide, "Algorithm 3 — LAD TreeBoost")
    add_accent_underline(slide)

    add_textbox(slide, Inches(0.55), Inches(1.45), Inches(1.5), Inches(0.5),
                text="With", size=18, color=TEXT, anchor=MSO_ANCHOR.MIDDLE)
    add_math_image(
        slide, Inches(1.10), Inches(1.45), Inches(4.5), Inches(0.5),
        latex=r"L(y, F) \;=\; |\,y - F\,|\, :",
        fontsize=22, align="left",
    )

    rows = [
        ("Pseudo-response:",
         r"\tilde{y}_i \;=\; \mathrm{sign}\!\left(y_i - F_{m-1}(\mathbf{x}_i)\right).",
         True),
        ("Tree fit:",
         "by least squares to the signs of the residuals.",
         False),
        ("Initial constant:",
         r"F_0 \;=\; \mathrm{median}(y).",
         True),
    ]
    y = Inches(2.10)
    for label, content, is_math in rows:
        add_textbox(slide, Inches(0.9), y, Inches(0.3), Inches(0.5),
                    text="•", size=18, bold=True, color=ACCENT,
                    anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(1.2), y, Inches(2.7), Inches(0.5),
                    text=label, size=16, bold=True, color=ACCENT,
                    anchor=MSO_ANCHOR.MIDDLE)
        if is_math:
            add_math_image(
                slide, Inches(3.95), y, Inches(8.7), Inches(0.5),
                latex=content, fontsize=18, align="left",
            )
        else:
            add_textbox(
                slide, Inches(3.95), y, Inches(8.7), Inches(0.5),
                text=content, size=15, color=TEXT,
                anchor=MSO_ANCHOR.MIDDLE,
            )
        y += Inches(0.55)

    add_keypoint_box(
        slide, Inches(0.55), Inches(3.85), Inches(12.2), Inches(1.10),
        "Crucial step: the per-leaf update is the median of residuals in the "
        "leaf, not the leaf mean. The tree's leaf values are overwritten with "
        "these medians.",
        size=16,
    )

    add_keypoint_math(
        slide, Inches(0.95), Inches(5.20), Inches(11.4), Inches(1.05),
        latex=r"\gamma_{jm} \;=\; \mathrm{median}_{\,\mathbf{x}_i \in R_{jm}}"
              r"\!\left(\, y_i - F_{m-1}(\mathbf{x}_i)\, \right)",
        fontsize=30,
    )

    add_textbox(
        slide, Inches(0.55), Inches(6.40), Inches(12.2), Inches(0.4),
        text=("This is equation (18) specialized to L₁. It is what makes the "
              "method robust."),
        size=12, italic=True, color=MUTED,
    )
    add_footer(slide, 9, TOTAL_SLIDES)


# ---------- Slide 10 — Huber ---------------------------------------------
def slide_10_huber():
    slide = add_blank_slide(prs)
    add_title(slide, "Algorithm 4 — M TreeBoost (Huber)")
    add_accent_underline(slide)

    add_textbox(slide, Inches(0.55), Inches(1.45), Inches(12.5), Inches(0.45),
                text="Quadratic for small residuals, linear in the tail:",
                size=17, color=TEXT)

    # Piecewise Huber loss: render as a card with "L(y, F) =", a tall brace
    # text shape, and two case lines stacked to the right of the brace.
    huber_x = Inches(0.55)
    huber_y = Inches(2.00)
    huber_w = Inches(12.2)
    huber_h = Inches(1.50)
    add_rect(slide, huber_x, huber_y, huber_w, huber_h,
             fill_color=CARD_BG, line_color=BORDER)
    # LHS  "L(y, F) ="
    add_math_image(
        slide, huber_x + Inches(0.40), huber_y + Inches(0.30),
        Inches(2.2), Inches(0.90),
        latex=r"L(y, F) \;=", fontsize=24, align="left",
    )
    # Big brace as a regular text shape (Cambria scales nicely).
    brace_tb = slide.shapes.add_textbox(
        huber_x + Inches(2.85), huber_y + Inches(0.05),
        Inches(0.55), Inches(1.40),
    )
    btf = brace_tb.text_frame
    btf.margin_left = Emu(0)
    btf.margin_right = Emu(0)
    btf.margin_top = Emu(0)
    btf.margin_bottom = Emu(0)
    btf.vertical_anchor = MSO_ANCHOR.MIDDLE
    btf.word_wrap = False
    bp = btf.paragraphs[0]
    bp.alignment = PP_ALIGN.CENTER
    brun = bp.add_run()
    brun.text = "{"
    brun.font.name = "Cambria"
    brun.font.size = Pt(80)
    brun.font.color.rgb = TEXT
    # The two case branches.
    add_math_image(
        slide, huber_x + Inches(3.40), huber_y + Inches(0.10),
        Inches(8.6), Inches(0.55),
        latex=r"\frac{1}{2}\,(y - F)^{2}"
              r"\quad,\qquad |\,y - F\,| \leq \delta",
        fontsize=22, align="left",
    )
    add_math_image(
        slide, huber_x + Inches(3.40), huber_y + Inches(0.80),
        Inches(8.6), Inches(0.55),
        latex=r"\delta\,|\,y - F\,| \;-\; \frac{1}{2}\,\delta^{2}"
              r"\quad,\qquad |\,y - F\,| > \delta",
        fontsize=22, align="left",
    )

    # Bullets — each contains a label + small math image.
    rows = [
        ("Adaptive δ:",
         r"\delta_m \;=\; \mathrm{quantile}_{\alpha}\!\left(|\,r_i\,|\right)"
         r"\quad\text{at each iteration.}"),
        ("Pseudo-response:",
         r"\tilde{y}_i \;=\; r_i\;\;\text{if}\;\; |r_i| \leq \delta_m,"
         r"\quad\text{else}\quad \delta_m\,\mathrm{sign}(r_i)."),
    ]
    y = Inches(3.65)
    for label, latex in rows:
        add_textbox(slide, Inches(0.9), y, Inches(0.3), Inches(0.5),
                    text="•", size=18, bold=True, color=ACCENT,
                    anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(1.2), y, Inches(2.5), Inches(0.5),
                    text=label, size=15, bold=True, color=ACCENT,
                    anchor=MSO_ANCHOR.MIDDLE)
        add_math_image(
            slide, Inches(3.75), y, Inches(8.9), Inches(0.5),
            latex=latex, fontsize=16, align="left",
        )
        y += Inches(0.45)

    # Sentence introducing the leaf-update formula (regular text, not math).
    add_textbox(
        slide, Inches(0.55), y + Inches(0.05), Inches(12.2), Inches(0.4),
        text=("Leaf update is Friedman's 1-step Newton-like adjustment "
              "around the leaf median:"),
        size=15, color=TEXT, anchor=MSO_ANCHOR.MIDDLE,
    )

    # Big leaf-update formula.
    add_keypoint_math(
        slide, Inches(0.55), Inches(5.20), Inches(12.2), Inches(1.10),
        latex=r"\gamma_{jm} \;=\; \tilde{r}_{jm} \,+\, \frac{1}{N_{jm}}"
              r"\sum_{\mathbf{x}_i \in R_{jm}} \mathrm{sign}(r_i - \tilde{r}_{jm})"
              r"\,\cdot\,\min\!\left(\delta_m,\; |\,r_i - \tilde{r}_{jm}\,|\right)",
        fontsize=22,
    )
    # Footnote: leaf-residual median, in math.
    add_textbox(slide, Inches(0.55), Inches(6.42), Inches(2.4), Inches(0.4),
                text="where", size=12, italic=True, color=MUTED,
                anchor=MSO_ANCHOR.MIDDLE)
    add_math_image(
        slide, Inches(1.20), Inches(6.42), Inches(11.0), Inches(0.4),
        latex=r"\tilde{r}_{jm} \;=\; \mathrm{median}_{\,\mathbf{x}_i \in R_{jm}}"
              r"\!\left(\, r_{m-1}(\mathbf{x}_i)\, \right).",
        fontsize=14, align="left",
    )
    add_footer(slide, 10, TOTAL_SLIDES)


# ---------- Slide 11 — Terminal-region update ----------------------------
def slide_11_terminal():
    slide = add_blank_slide(prs)
    add_title(slide, "The terminal-region update — the secret sauce")
    add_accent_underline(slide)

    add_textbox(slide, Inches(0.55), Inches(1.45), Inches(12.2), Inches(0.4),
                text="Two stages per iteration when the weak learner is a tree:",
                size=16, color=TEXT)

    steps = [
        ("Build", "a J-leaf tree by least squares on the pseudo-responses ỹ."),
        ("Re-solve", "a tiny 1-D problem inside each leaf using the original loss:"),
    ]
    y = Inches(2.0)
    for label, body in steps:
        add_textbox(slide, Inches(0.85), y, Inches(0.5), Inches(0.4),
                    text="1." if label == "Build" else "2.",
                    size=17, bold=True, color=ACCENT)
        tb = slide.shapes.add_textbox(Inches(1.3), y, Inches(11.5), Inches(0.5))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = label
        run.font.bold = True
        run.font.size = Pt(16)
        run.font.color.rgb = ACCENT
        run.font.name = FONT_BODY
        run = p.add_run()
        run.text = " " + body
        run.font.size = Pt(16)
        run.font.color.rgb = TEXT
        run.font.name = FONT_BODY
        y += Inches(0.55)

    add_keypoint_math(
        slide, Inches(1.3), y, Inches(11.0), Inches(1.00),
        latex=r"\gamma_{jm} \;=\; \arg\min_{\gamma}\,"
              r"\sum_{\mathbf{x}_i \in R_{jm}}\, L\!\left(y_i,\; "
              r"F_{m-1}(\mathbf{x}_i) + \gamma\right)",
        fontsize=24,
    )
    y += Inches(1.20)
    add_textbox(slide, Inches(0.85), y, Inches(0.5), Inches(0.4),
                text="3.", size=17, bold=True, color=ACCENT,
                anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, Inches(1.3), y, Inches(1.4), Inches(0.4),
                text="Update:", size=16, bold=True, color=ACCENT,
                anchor=MSO_ANCHOR.MIDDLE)
    add_math_image(
        slide, Inches(2.70), y, Inches(10.0), Inches(0.4),
        latex=r"F_m(\mathbf{x}) \;=\; F_{m-1}(\mathbf{x}) + \nu\,\gamma_{jm},"
              r"\quad \mathbf{x} \in R_{jm}.",
        fontsize=18, align="left",
    )

    add_keypoint_box(
        slide, Inches(0.55), Inches(5.85), Inches(12.2), Inches(1.15),
        "Why this is the crux:  the LS-fit tree is a fast, smooth gradient "
        "direction; the loss-specific γ_jm re-injects the original loss's "
        "geometry — robust medians for LAD, clipped sums for Huber, Newton "
        "steps for logistic.",
        size=14,
    )
    add_footer(slide, 11, TOTAL_SLIDES)


# ---------- Slide 12 — Implementation architecture -----------------------
def slide_12_impl():
    slide = add_blank_slide(prs)
    add_title(slide, "From-scratch implementation")
    add_accent_underline(slide)

    add_textbox(slide, Inches(0.55), Inches(1.45), Inches(12.2), Inches(0.4),
                text="One clean abstraction does all the work:",
                size=16, color=TEXT)
    code = (
        "class Loss(ABC):\n"
        "    def initial_prediction(self, y):           ...   # F_0\n"
        "    def negative_gradient(self, y, F):         ...   # ỹ_i   (line 3)\n"
        "    def leaf_update(self, y, F, indices):      ...   # γ_jm  (eq. 18)\n"
        "    def update_state(self, y, F):              ...   # e.g. Huber's δ_m"
    )
    add_code_block(slide, Inches(0.55), Inches(2.00),
                   Inches(12.2), Inches(1.85), code, size=14)

    add_textbox(
        slide, Inches(0.55), Inches(3.95), Inches(12.2), Inches(0.4),
        text=("Three losses implement this contract; the boosting loop is "
              "loss-agnostic."),
        size=14, italic=True, color=MUTED,
    )

    cards = [
        ("tree.py", "CART-style regression tree, best-first growth, vectorized "
                    "split search. set_leaf_values() lets the boosting driver "
                    "overwrite leaf means with γ."),
        ("losses.py", "LS / LAD / Huber, each ~30 LOC. HuberLoss.update_state "
                      "sets δ_m at each iteration."),
        ("model.py", "TreeBoostRegressor drives Algorithms 2 / 3 / 4 with "
                     "shrinkage ν."),
    ]
    card_w = Inches(4.0)
    card_h = Inches(2.05)
    gap = Inches(0.15)
    x = Inches(0.55)
    for name, body in cards:
        add_rect(slide, x, Inches(4.55), card_w, card_h,
                 fill_color=WHITE, line_color=BORDER)
        add_pill(slide, x + Inches(0.25), Inches(4.7), name)
        tb = slide.shapes.add_textbox(x + Inches(0.25), Inches(5.10),
                                      card_w - Inches(0.5), card_h - Inches(0.6))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = body
        run.font.name = FONT_BODY
        run.font.size = Pt(13)
        run.font.color.rgb = TEXT
        x += card_w + gap

    add_textbox(
        slide, Inches(0.55), Inches(6.85), Inches(12.2), Inches(0.4),
        text=("No scikit-learn, XGBoost, LightGBM, or CatBoost. "
              "Only numpy for arrays. 30 passing pytest tests."),
        size=12, italic=True, color=MUTED,
    )
    add_footer(slide, 12, TOTAL_SLIDES)


# ---------- Slide 13 — Custom regression tree ----------------------------
def slide_13_tree():
    slide = add_blank_slide(prs)
    add_title(slide, "Custom regression tree (the weak learner)")
    add_accent_underline(slide)

    bullets = [
        ("Best-first growth",
         " — repeatedly expand the leaf with the highest SSE-reduction split "
         "until max_leaves = J is reached."),
        ("Vectorized split search",
         " — sort each feature column, pre-compute cumulative sums of ỹ, "
         "evaluate all valid split points in one shot:"),
    ]
    y = Inches(1.55)
    for label, body in bullets:
        add_textbox(slide, Inches(0.85), y, Inches(0.3), Inches(0.4),
                    text="•", size=18, bold=True, color=ACCENT)
        tb = slide.shapes.add_textbox(Inches(1.2), y, Inches(11.5), Inches(0.95))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = label
        run.font.bold = True
        run.font.size = Pt(15)
        run.font.color.rgb = ACCENT
        run.font.name = FONT_BODY
        run = p.add_run()
        run.text = body
        run.font.size = Pt(15)
        run.font.color.rgb = TEXT
        run.font.name = FONT_BODY
        y += Inches(0.95)

    add_keypoint_math(
        slide, Inches(2.5), Inches(3.7), Inches(8.3), Inches(1.00),
        latex=r"\mathrm{gain}(s) \;=\; \frac{S_{L}^{\,2}}{n_{L}}"
              r" \;+\; \frac{S_{R}^{\,2}}{n_{R}}"
              r" \;-\; \frac{S^{\,2}}{n}",
        fontsize=28,
    )

    extra = [
        ("Honors J directly",
         " — same complexity knob the paper uses (Section 5)."),
        ("Leaf-id API (apply, set_leaf_values)",
         " — boosting driver maps each row to its leaf, computes γ_jm per "
         "leaf, and overwrites the LS leaf means."),
    ]
    y = Inches(4.85)
    for label, body in extra:
        add_textbox(slide, Inches(0.85), y, Inches(0.3), Inches(0.4),
                    text="•", size=18, bold=True, color=ACCENT)
        tb = slide.shapes.add_textbox(Inches(1.2), y, Inches(11.5), Inches(0.85))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = label
        run.font.bold = True
        run.font.size = Pt(15)
        run.font.color.rgb = ACCENT
        run.font.name = FONT_BODY
        run = p.add_run()
        run.text = body
        run.font.size = Pt(15)
        run.font.color.rgb = TEXT
        run.font.name = FONT_BODY
        y += Inches(0.75)

    add_textbox(
        slide, Inches(0.55), Inches(6.6), Inches(12), Inches(0.4),
        text="Roughly 200 lines, deterministic, and readable. Built for correctness over throughput.",
        size=12, italic=True, color=MUTED,
    )
    add_footer(slide, 13, TOTAL_SLIDES)


# ---------- Slide 14 — Experiment 1 --------------------------------------
def slide_14_exp1():
    slide = add_blank_slide(prs)
    add_title(slide, "Experiment 1 — convergence on a clean signal")
    add_accent_underline(slide)

    # Left column
    add_textbox(slide, Inches(0.55), Inches(1.45), Inches(6), Inches(0.4),
                text="Friedman-#1 target", size=15, bold=True, color=ACCENT)
    add_math_image(
        slide, Inches(0.55), Inches(1.80), Inches(6.2), Inches(0.55),
        latex=r"f(\mathbf{x}) \,=\, 10\sin(\pi x_1 x_2)"
              r" + 20(x_3 - 0.5)^{2} + 10\,x_4 + 5\,x_5",
        fontsize=15, align="left",
    )
    add_textbox(slide, Inches(0.55), Inches(2.40), Inches(1.9), Inches(0.4),
                text="Train / test:", size=13, bold=True, color=ACCENT,
                anchor=MSO_ANCHOR.MIDDLE)
    add_math_image(
        slide, Inches(2.45), Inches(2.40), Inches(4.3), Inches(0.4),
        latex=r"n_{\mathrm{train}} = n_{\mathrm{test}} = 600,"
              r"\quad \sigma = 1.0\;\;\text{(Gaussian noise).}",
        fontsize=13, align="left",
    )
    add_textbox(slide, Inches(0.55), Inches(2.80), Inches(1.9), Inches(0.4),
                text="Boosting:", size=13, bold=True, color=ACCENT,
                anchor=MSO_ANCHOR.MIDDLE)
    add_math_image(
        slide, Inches(2.45), Inches(2.80), Inches(4.3), Inches(0.4),
        latex=r"M = 200,\quad \nu = 0.1,\quad J = 8"
              r"\;\;\text{for all three losses.}",
        fontsize=13, align="left",
    )

    table_data = [
        ["Loss", "test MSE", "test MAE"],
        ["LS", "1.989", "1.136"],
        ["LAD", "2.309", "1.194"],
        ["Huber", "1.977", "1.124"],
    ]
    add_table(slide, Inches(0.55), Inches(3.30), Inches(6.0), Inches(1.85),
              table_data, num_cols_right={1, 2}, font_size=14)

    add_keypoint_box(
        slide, Inches(0.55), Inches(5.30), Inches(6.0), Inches(1.30),
        "Under Gaussian noise: LS ≈ Huber; LAD pays the classical small "
        "efficiency cost.",
        size=14,
    )

    # Right column
    add_picture_card(slide, Inches(6.85), Inches(1.45),
                     Inches(6.0), Inches(4.4),
                     RESULTS / "convergence_clean.png")
    add_textbox(slide, Inches(6.85), Inches(5.95), Inches(6.0), Inches(0.4),
                text="Train (blue) and validation (orange) loss curves per iteration.",
                size=11, italic=True, color=MUTED)
    add_footer(slide, 14, TOTAL_SLIDES)


# ---------- Slide 15 — Experiment 2 --------------------------------------
def slide_15_exp2():
    slide = add_blank_slide(prs)
    add_title(slide, "Experiment 2 — robustness under contamination")
    add_accent_underline(slide)

    add_textbox(
        slide, Inches(0.55), Inches(1.45), Inches(6.0), Inches(0.55),
        text="Same target, but 10% of training labels perturbed by "
             "heavy-tailed shocks",
        size=14, color=TEXT,
    )
    add_math_image(
        slide, Inches(0.55), Inches(2.05), Inches(6.0), Inches(0.45),
        latex=r"y_i \;\leftarrow\; y_i \,\pm\, 25\,t_{3}"
              r"\quad\text{(test set is clean).}",
        fontsize=16, align="left",
    )

    table_data = [
        ["Loss", "clean MSE", "clean MAE"],
        ["LS", "33.04", "3.63"],
        ["LAD", "3.49", "1.17"],
        ["Huber", "11.28", "1.84"],
    ]
    add_table(slide, Inches(0.55), Inches(2.65), Inches(6.0), Inches(1.85),
              table_data, num_cols_right={1, 2}, font_size=14)

    bullets = [
        "LS collapses — squared loss is dominated by outliers.",
        "LAD recovers almost completely (medians ignore extremes).",
        "Huber is much better than LS, ranking depends on δ.",
    ]
    y = Inches(4.65)
    for item in bullets:
        add_textbox(slide, Inches(0.55), y, Inches(0.3), Inches(0.4),
                    text="•", size=16, bold=True, color=ACCENT)
        add_textbox(slide, Inches(0.85), y, Inches(5.7), Inches(0.6),
                    text=item, size=13, color=TEXT)
        y += Inches(0.6)

    add_picture_card(slide, Inches(6.85), Inches(1.45),
                     Inches(6.0), Inches(4.4),
                     RESULTS / "robustness_bar.png")
    add_textbox(slide, Inches(6.85), Inches(5.95), Inches(6.0), Inches(0.4),
                text="Clean-test error per loss under 10% contamination.",
                size=11, italic=True, color=MUTED)
    add_footer(slide, 15, TOTAL_SLIDES)


# ---------- Slide 16 — Predicted vs true ---------------------------------
def slide_16_scatter():
    slide = add_blank_slide(prs)
    add_title(slide, "Predicted vs true under contamination")
    add_accent_underline(slide)

    add_picture_card(slide, Inches(1.0), Inches(1.45),
                     Inches(11.3), Inches(4.6),
                     RESULTS / "robustness_true_vs_pred.png")
    add_keypoint_box(
        slide, Inches(0.55), Inches(6.20), Inches(12.2), Inches(0.85),
        "LS predictions drift well off the diagonal; LAD and Huber stay "
        "calibrated despite the noisy training labels. This is the picture "
        "Section 4 of the paper promises.",
        size=14,
    )
    add_footer(slide, 16, TOTAL_SLIDES)


# ---------- Slide 17 — Shrinkage -----------------------------------------
def slide_17_shrinkage():
    slide = add_blank_slide(prs)
    add_title(slide, "Experiment 3 — shrinkage trade-off")
    add_accent_underline(slide)

    add_textbox(slide, Inches(0.55), Inches(1.45), Inches(6.0), Inches(0.4),
                text="LS Boost on the clean signal, sweeping the learning rate ν.",
                size=14, color=TEXT)

    bullets = [
        ("ν = 1.0", " : hits the floor in ~30 iterations, then plateaus."),
        ("ν = 0.1", " : lower final validation MSE — but takes ~150–200 trees."),
        ("ν = 0.05", " : slowest; needs more iterations to compete."),
    ]
    y = Inches(2.05)
    for label, body in bullets:
        add_textbox(slide, Inches(0.55), y, Inches(0.3), Inches(0.4),
                    text="•", size=16, bold=True, color=ACCENT)
        tb = slide.shapes.add_textbox(Inches(0.85), y, Inches(5.7), Inches(0.6))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = label
        run.font.bold = True
        run.font.size = Pt(14)
        run.font.color.rgb = ACCENT
        run.font.name = FONT_BODY
        run = p.add_run()
        run.text = body
        run.font.size = Pt(14)
        run.font.color.rgb = TEXT
        run.font.name = FONT_BODY
        y += Inches(0.6)

    add_keypoint_box(
        slide, Inches(0.55), Inches(4.0), Inches(6.0), Inches(1.40),
        "Reproduces the standard recommendation that emerged from this paper: "
        "moderate shrinkage + larger M generalizes better.",
        size=14,
    )

    add_picture_card(slide, Inches(6.85), Inches(1.45),
                     Inches(6.0), Inches(4.4),
                     RESULTS / "shrinkage_curves.png")
    add_textbox(slide, Inches(6.85), Inches(5.95), Inches(6.0), Inches(0.4),
                text="Validation MSE/2 vs boosting iteration for several ν values.",
                size=11, italic=True, color=MUTED)
    add_footer(slide, 17, TOTAL_SLIDES)


# ---------- Slide 18 — What worked / didn't ------------------------------
def slide_18_worked():
    slide = add_blank_slide(prs)
    add_title(slide, "What worked & what didn't")
    add_accent_underline(slide)

    # Left card — Worked
    add_rect(slide, Inches(0.55), Inches(1.45), Inches(6.05), Inches(5.3),
             fill_color=WHITE, line_color=BORDER)
    add_pill(slide, Inches(0.75), Inches(1.65), "Worked",
             fg=GOOD_TEXT, bg=GOOD_BG)
    worked = [
        "The 3-method Loss interface kept the boosting loop tiny and "
        "loss-agnostic.",
        "Best-first tree growth honored J cleanly with no extra plumbing.",
        "Unit tests written against the paper's equations (e.g. Algorithm 4 "
        "leaf update) caught the real bugs early.",
        "Experiments reproduced the qualitative paper claims: LS = reality "
        "check, LAD/Huber robust, shrinkage helps.",
    ]
    y = Inches(2.10)
    for item in worked:
        add_textbox(slide, Inches(0.85), y, Inches(0.3), Inches(0.4),
                    text="•", size=16, bold=True, color=GOOD_TEXT)
        add_textbox(slide, Inches(1.15), y, Inches(5.4), Inches(1.2),
                    text=item, size=12, color=TEXT)
        y += Inches(1.0)

    # Right card — Limits
    add_rect(slide, Inches(6.75), Inches(1.45), Inches(6.05), Inches(5.3),
             fill_color=WHITE, line_color=BORDER)
    add_pill(slide, Inches(6.95), Inches(1.65), "Honest limits",
             fg=WARN_TEXT, bg=WARN_BG)
    limits = [
        "CART splitter is correct but slow — fine for a few hundred trees, "
        "not industrial.",
        "No subsampling, no influence trimming (regression scope by design).",
        "Huber's exact ranking depends strongly on the quantile chosen for δ.",
        "LAD training loss is non-monotone per-iteration because sign(r) "
        "discards magnitude.",
    ]
    y = Inches(2.10)
    for item in limits:
        add_textbox(slide, Inches(7.05), y, Inches(0.3), Inches(0.4),
                    text="•", size=16, bold=True, color=WARN_TEXT)
        add_textbox(slide, Inches(7.35), y, Inches(5.4), Inches(1.2),
                    text=item, size=12, color=TEXT)
        y += Inches(1.0)
    add_footer(slide, 18, TOTAL_SLIDES)


# ---------- Slide 19 — Insights & legacy ---------------------------------
def slide_19_legacy():
    slide = add_blank_slide(prs)
    add_title(slide, "What the results reveal & lasting impact")
    add_accent_underline(slide)

    add_textbox(slide, Inches(0.55), Inches(1.45), Inches(12.2), Inches(0.4),
                text="What the results reveal", size=18, bold=True, color=ACCENT)
    insights = [
        "Function-space framing isn't just exposition. Once Loss exposes "
        "three methods, adding a new regression loss takes minutes.",
        "The terminal-region update is the practical heart. Without it, LAD "
        "would degenerate; with it, leaves carry medians and the method is "
        "genuinely robust.",
        "Robustness has a small efficiency cost. LAD ≈16% worse than LS on "
        "clean data; LS ≈10× worse than LAD under contamination — the "
        "asymmetry favors using a robust default.",
        "Shrinkage = nearly-free regularization. Smaller ν, larger M, better "
        "generalization, more compute.",
    ]
    y = Inches(1.95)
    for item in insights:
        add_textbox(slide, Inches(0.85), y, Inches(0.3), Inches(0.4),
                    text="•", size=14, bold=True, color=ACCENT)
        add_textbox(slide, Inches(1.15), y, Inches(11.5), Inches(0.7),
                    text=item, size=13, color=TEXT)
        y += Inches(0.65)

    add_textbox(slide, Inches(0.55), Inches(4.85), Inches(12.2), Inches(0.4),
                text="The paper's lasting impact", size=18, bold=True, color=ACCENT)
    legacy = [
        "Unified boosting. Loss-agnostic boosting became the standard framing.",
        "Modern descendants: XGBoost (2014, 2nd-order, regularization, "
        "sparsity), LightGBM (2017, histograms, leaf-wise growth), CatBoost "
        "(2018, ordered boosting, target encoding).",
        "Interpretation tools (variable importance, partial dependence) come "
        "straight from this paper's TreeBoost section.",
    ]
    y = Inches(5.30)
    for item in legacy:
        add_textbox(slide, Inches(0.85), y, Inches(0.3), Inches(0.4),
                    text="•", size=14, bold=True, color=ACCENT)
        add_textbox(slide, Inches(1.15), y, Inches(11.5), Inches(0.7),
                    text=item, size=13, color=TEXT)
        y += Inches(0.55)
    add_footer(slide, 19, TOTAL_SLIDES)


# ---------- Slide 20 — Q&A -----------------------------------------------
def slide_20_qna():
    slide = add_blank_slide(prs)
    add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, fill_color=WHITE)
    add_rect(slide, 0, Inches(2.4), prs.slide_width, Inches(0.06),
             fill_color=ACCENT)

    add_textbox(slide, Inches(0.7), Inches(2.7), Inches(12), Inches(1.0),
                text="Thank you", size=54, bold=True, color=ACCENT)
    add_textbox(slide, Inches(0.7), Inches(3.85), Inches(12), Inches(0.6),
                text="Questions?", size=28, italic=True, color=MUTED)

    add_textbox(
        slide, Inches(0.7), Inches(5.0), Inches(12), Inches(0.5),
        text=("Code, tests, experiments, and analysis writeup are all in the "
              "project repo."),
        size=15, color=TEXT,
    )
    add_textbox(
        slide, Inches(0.7), Inches(5.5), Inches(12), Inches(0.5),
        text="src/treeboost/   ·   tests/   ·   experiments/   ·   analysis.md   ·   README.md",
        size=13, italic=True, color=MUTED, font=FONT_MONO,
    )


# ---------------------------------------------------------------------------
# Run
# ---------------------------------------------------------------------------

builders = [
    slide_1_title, slide_2_roadmap, slide_3_problem, slide_4_literature,
    slide_5_additive, slide_6_insight, slide_7_algo1, slide_8_ls,
    slide_9_lad, slide_10_huber, slide_11_terminal, slide_12_impl,
    slide_13_tree, slide_14_exp1, slide_15_exp2, slide_16_scatter,
    slide_17_shrinkage, slide_18_worked, slide_19_legacy, slide_20_qna,
]
assert len(builders) == TOTAL_SLIDES, "slide count mismatch"

for builder in builders:
    builder()

prs.save(OUTPUT)
print(f"Wrote {OUTPUT.relative_to(ROOT)}  ({len(prs.slides)} slides)")
